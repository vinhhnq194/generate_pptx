from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from .theme import BORDER, SURFACE, LABEL

def add_placeholder(slide, left, top, width, height, label=None, name=None):
    """
    Create a rounded-rectangle placeholder with optional text label and .name set.
    """
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = SURFACE
    shp.line.color.rgb = BORDER
    try:
        shp.adjustments[0] = 0.16
    except Exception:
        pass

    # assign a programmatic name (for later reference)
    if name:
        shp.name = name
    elif label:
        shp.name = f"Placeholder-{label}"
    else:
        shp.name = "Placeholder"

    # optional visible label inside the shape
    if label:
        tb = slide.shapes.add_textbox(
            Inches(left + 0.08), Inches(top + 0.06),
            Inches(width - 0.16), Inches(0.32)
        )
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(9)
        p.font.color.rgb = LABEL

    return shp
