import json
from pptx import Presentation
from pptx.util import Inches
from v6.parsers.layout_parser import parse_layout_tree
from v6.layout.grid_solver import solve_layout_tree
from v6.layout.placement_debug import draw_grid, draw_bbox
from v6.renderer.primitives import add_placeholder

def render_layout_only(html_path: str, styles_path: str, template_path: str | None = None):
    ST = json.load(open(styles_path, "r", encoding="utf-8"))
    page, bands = ST["page"], ST["bands"]

    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width  = Inches(page["width_in"])
    prs.slide_height = Inches(page["height_in"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tree = parse_layout_tree(html_path)
    placements = solve_layout_tree(tree, page, bands)

    if ST.get("debug", {}).get("grid", False):
        draw_grid(slide, page, page["gutter_in"])

    for i, (rect, col) in enumerate(placements, 1):
        label = f"col:{col.span} off:{col.offset}"
        name  = f"Col_span{col.span}_off{col.offset}"
        add_placeholder(slide, rect.left, rect.top, rect.width, rect.height, label=label, name=name)
        if ST.get("debug", {}).get("bbox", False):
            draw_bbox(slide, rect)

    return prs
