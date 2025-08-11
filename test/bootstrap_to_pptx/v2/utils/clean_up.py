from pptx.shapes.base import BaseShape

SAFE_KEEP_NAMES = {
    "DECOR_DIAGONAL", "DECOR_CIRCLE",
    "CARD_NARRATIVE", "CARD_STEPS", "CARD_ICONS", "CARD_OUTLOOK",
    "KPI_TILE", "FOOTER_BAR"
}

def _is_textbox_empty(shape: BaseShape) -> bool:
    try:
        if not hasattr(shape, "text_frame"): return False
        txt = shape.text_frame.text if shape.text_frame else ""
        return (txt or "").strip() == ""
    except Exception:
        return False

def delete_shape(shape: BaseShape):
    # Official API lacks delete; this is the accepted approach.
    sp = shape._element
    sp.getparent().remove(sp)

def cleanup_slide(slide):
    """
    Remove stray/empty textboxes created by guards or autosize edge cases.
    Keeps any shape we have explicitly named in SAFE_KEEP_NAMES.
    """
    # Collect first, then remove (can't modify while iterating)
    to_delete = []
    for shp in slide.shapes:
        nm = (getattr(shp, "name", "") or "").upper()
        if nm in SAFE_KEEP_NAMES:
            continue
        # delete empty textboxes or 0-size pictures
        if _is_textbox_empty(shp):
            to_delete.append(shp)
            continue
        if shp.height == 0 or shp.width == 0:
            to_delete.append(shp)

    for shp in to_delete:
        delete_shape(shp)
