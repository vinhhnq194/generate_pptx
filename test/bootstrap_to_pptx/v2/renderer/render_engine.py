import os
import json
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

from parsers.bootstrap_html_to_ilt import parse_bootstrap_html_to_ilt
from renderer.layout_solver import layout_rows
from renderer.element_factory import build_column_contents
from utils.merge import deep_update
from utils.clean_up import cleanup_slide

from .grid import Grid12
from .elements import (
    add_title, add_subtitle, add_decor_diagonal, add_decor_circle,
    add_card, add_card_header, add_text, add_text_padded, add_bullets,
    add_kpi_tile, add_footer_bar
)
from schema.slide_model import SlideModel
from utils.text_fit import fit_font_size, wrap_text

# Configure logging for this module; adjust level in your app entrypoint if needed
logging.basicConfig(level=logging.DEBUG, format="%(levelname)s: %(message)s")

def _load_json(path):
    logging.debug(f"Loading JSON file: {path}")
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    logging.debug(f"Loaded JSON keys: {list(data)[:10]}{'...' if len(data) > 10 else ''}")
    return data

def _has_text(s):
    """Truthy if string has non-space or object has any non-empty text-like field."""
    if s is None:
        return False
    if isinstance(s, str):
        result = bool(s.strip())
        logging.debug(f"_has_text(str): {result} (sample='{s[:30]}...')")
        return result
    for field in ("text", "title", "subtitle", "header", "caption", "headline"):
        v = getattr(s, field, None)
        if isinstance(v, str) and v.strip():
            logging.debug(f"_has_text(obj): True via field '{field}' (sample='{v[:30]}...')")
            return True
    logging.debug("_has_text(obj): False")
    return False

def _has_list(lst):
    result = bool(lst and any(_has_text(x) for x in lst))
    logging.debug(f"_has_list: {result} (len={len(lst) if lst else 0})")
    return result

def render_from_html(html_path: str, mapping_path: str, styles: dict, overrides: dict, template_path: str = None):
    logging.info("=== render_from_html start ===")
    logging.debug(f"Args: html_path={html_path}, mapping_path={mapping_path}, template_path={template_path}")

    # load mapping.json
    mapping = _load_json(mapping_path)

    # build ILT
    logging.info("Parsing HTML -> ILT (intermediate layout tree)")
    ilt = parse_bootstrap_html_to_ilt(html_path, mapping)
    logging.debug(f"ILT: decor={ilt.decor}, title={ilt.title}, subtitle={ilt.subtitle}, rows={len(ilt.rows)}")

    # merge styles + overrides
    logging.info("Merging styles with overrides")
    ST = deep_update(styles, overrides or {})
    logging.debug(f"Merged style keys: {list(ST.keys())}")

    # Presentation
    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width  = Inches(ST["page"]["width_in"])
    prs.slide_height = Inches(ST["page"]["height_in"])
    logging.debug(f"Presentation size set to {ST['page']['width_in']}x{ST['page']['height_in']} inches")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logging.debug("Blank slide added")

    # Grid
    grid = Grid12(
        slide_width_in=ST["page"]["width_in"],
        slide_height_in=ST["page"]["height_in"],
        margins_in=tuple(ST["page"]["margins_in"]),
        gutter_in=ST["page"]["gutter_in"]
    )
    logging.debug("Grid12 initialized")

    # Decor, title, subtitle
    if "diagonal" in ilt.decor:
        logging.debug("Adding decor: diagonal")
        add_decor_diagonal(slide, -2.5, -2.0, 16.0, 10.0, -10, (13,110,253,0.85))
    if "circle" in ilt.decor:
        logging.debug("Adding decor: circle")
        add_decor_circle(slide, 11.0, 6.4, 2.5, (13,110,253,0.85))
    if ilt.title:
        logging.debug(f"Adding title: '{ilt.title[:60]}...'")
        add_title(slide, 0.6, ST["title"]["top_in"], 12.0, 0.9, ilt.title)
    if ilt.subtitle:
        logging.debug(f"Adding subtitle: '{ilt.subtitle[:60]}...'")
        add_subtitle(slide, 0.6, ST["subtitle"]["top_in"], 12.0, 0.6, ilt.subtitle)

    # Compute rows
    band_top = 1.8
    band_h   = 4.85
    logging.debug(f"Layout band: top={band_top} in, height={band_h} in")
    rows_layout = layout_rows(grid, ilt.rows, band_top, band_h, row_gap_in=0.0)
    logging.debug(f"Computed rows_layout: {sum(len(r) for r in rows_layout)} groups")

    # Render each column group
    logging.info("Rendering column groups")
    for row_idx, row in enumerate(rows_layout):
        for col_idx, (rect, group) in enumerate(row):
            logging.debug(f"Row {row_idx}, ColGroup {col_idx}: rect=({rect.left},{rect.top},{rect.width},{rect.height}), kind={getattr(group,'kind','column')}")
            build_column_contents(slide, rect, group, ST, {}, grid)

    # Footer
    if ilt.footer_left:
        logging.debug("Adding footer bar")
        add_footer_bar(slide, left=0.6, top=ST["page"]["height_in"] - (ST["footer"]["height_in"] + 0.35),
                       width=ST["page"]["width_in"] - 1.2, height=ST["footer"]["height_in"],
                       left_text=ilt.footer_left, right_text="Slide 1",
                       left_pt=ST["footer"]["left_pt"], right_pt=ST["footer"]["right_pt"])

    logging.debug("Cleaning slide")
    cleanup_slide(slide)
    logging.info("=== render_from_html end ===")
    return prs


def render_slide(model: SlideModel, *, template_path: str = None,
                 slide_num: int = 1, slide_count: int = 1,
                 overrides_path: str = None):
    logging.info("=== render_slide start ===")
    logging.debug(f"Args: template_path={template_path}, slide_num={slide_num}, slide_count={slide_count}, overrides_path={overrides_path}")

    here = os.path.dirname(__file__)
    cfg_dir = os.path.abspath(os.path.join(here, "..", "config"))
    styles_path = os.path.join(cfg_dir, "styles.json")
    icons_path  = os.path.join(cfg_dir, "icon_map.json")
    logging.debug(f"Config dir: {cfg_dir}")

    ST = _load_json(styles_path)

    ICON_MAP = {}
    try:
        ICON_MAP = _load_json(icons_path)
    except FileNotFoundError:
        logging.warning(f"Icon map not found: {icons_path}")
        ICON_MAP = {}

    # Apply overrides (optional)
    if overrides_path and os.path.exists(overrides_path):
        logging.info(f"Applying style overrides: {overrides_path}")
        OV = _load_json(overrides_path)
        deep_update(ST, OV)  # inplace merge
    else:
        logging.debug("No overrides applied")

    # Presentation
    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width  = Inches(ST["page"]["width_in"])
    prs.slide_height = Inches(ST["page"]["height_in"])
    logging.debug(f"Presentation size set to {ST['page']['width_in']}x{ST['page']['height_in']} inches")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    logging.debug("Blank slide added")

    # Grid
    grid = Grid12(
        slide_width_in = ST["page"]["width_in"],
        slide_height_in= ST["page"]["height_in"],
        margins_in     = tuple(ST["page"]["margins_in"]),
        gutter_in      = ST["page"]["gutter_in"]
    )
    logging.debug("Grid12 initialized")

    # Decor (optional)
    if (model.decor and ST.get("decor", {}).get("enabled", True)):
        logging.debug(f"Adding decor items: {len(model.decor)}")
        for d_idx, d in enumerate(model.decor):
            if d.kind == "diagonal":
                logging.debug(f"Decor {d_idx}: diagonal")
                shp = add_decor_diagonal(slide, left=-2.5, top=-2.0, width=16.0, height=10.0, angle=-10, rgba=(13,110,253,0.85))
                shp.name = "DECOR_DIAGONAL"
            elif d.kind == "circle":
                logging.debug(f"Decor {d_idx}: circle")
                shp = add_decor_circle(slide, left=11.0, top=6.4, diameter=2.5, rgba=(13,110,253,0.85))
                shp.name = "DECOR_CIRCLE"

    # Title / Subtitle (only if provided)
    if _has_text(getattr(model.title_block, "title", None)):
        logging.debug(f"Adding TITLE: '{model.title_block.title[:60]}...'")
        tbox = add_title(slide, 0.6, ST["title"]["top_in"], 12.0, 0.9, model.title_block.title)
        tbox.text_frame.paragraphs[0].font.size = Pt(ST["title"]["size_pt"])
        tbox.name = "TITLE"
    if _has_text(getattr(model.title_block, "subtitle", None)):
        logging.debug(f"Adding SUBTITLE: '{model.title_block.subtitle[:60]}...'")
        sbox = add_subtitle(slide, 0.6, ST["subtitle"]["top_in"], 12.0, 0.6, model.title_block.subtitle)
        sbox.text_frame.paragraphs[0].font.size = Pt(ST["subtitle"]["size_pt"])
        sbox.name = "SUBTITLE"

    # Bands (from style or overrides)
    row_top    = ST.get("bands", {}).get("row_top_in", 1.8)
    row_height = ST.get("bands", {}).get("row_height_in", 4.85)
    logging.debug(f"Bands: row_top={row_top}, row_height={row_height}")

    L_left, L_top, L_w, L_h = grid.rect_for(row_top_in=row_top, col_start=0, col_span=6, height_in=row_height)
    R_left, R_top, R_w, R_h = grid.rect_for(row_top_in=row_top, col_start=6, col_span=6, height_in=row_height)
    logging.debug(f"LEFT rect: {L_left},{L_top},{L_w},{L_h}")
    logging.debug(f"RIGHT rect: {R_left},{R_top},{R_w},{R_h}")

    # LEFT — create shapes only when there is content
    has_narr = bool(model.narrative and (_has_list(model.narrative.paragraphs) or _has_list(model.narrative.bullets)))
    has_kpi = bool(model.kpis and any(
        (_has_text(getattr(k, "headline", None))) or (_has_text(getattr(k, "caption", None)))
        for k in model.kpis
    ))
    logging.debug(f"LEFT: has_narr={has_narr}, has_kpi={has_kpi}")

    if has_narr or has_kpi:
        # Fix KPI band height; narrative uses remaining space
        kpi_h    = ST.get("left", {}).get("kpi_height_in", ST["kpi"]["height_in"])
        kpi_gap  = ST.get("left", {}).get("kpi_gap_in",     ST["kpi"]["gap_in"])
        card_h   = L_h - (kpi_h + 0.3) if has_kpi else L_h
        logging.debug(f"LEFT: kpi_h={kpi_h}, kpi_gap={kpi_gap}, narrative_card_h={card_h}")

        if has_narr:
            logging.debug("Adding narrative card")
            card = add_card(slide, L_left, L_top, L_w, card_h, radius=True, shadow=ST["shadow"]["enabled"])
            card.name = "CARD_NARRATIVE"

            pad_l, pad_t, pad_r, pad_b = ST["narrative"]["padding_in"]
            x = L_left + pad_l
            y = L_top  + pad_t
            w = L_w    - (pad_l + pad_r)
            h = card_h - (pad_t + pad_b)
            logging.debug(f"Narrative inner box: x={x}, y={y}, w={w}, h={h}")

            # budgets
            chars_per_line = int(35 * (w / 3.5))
            lines_budget   = int(h / 0.28)
            budget         = max(120, chars_per_line * lines_budget)
            logging.debug(f"Fit budgets: chars_per_line={chars_per_line}, lines={lines_budget}, budget={budget}")

            body_pt   = fit_font_size(" ".join(model.narrative.paragraphs or []), budget,
                                      base_pt=ST["narrative"]["body_size_pt"], min_pt=12)
            bullet_pt = fit_font_size(" ".join(model.narrative.bullets or []),   budget,
                                      base_pt=ST["narrative"]["bullets_size_pt"], min_pt=12)
            logging.debug(f"Fitted font sizes: body_pt={body_pt}, bullet_pt={bullet_pt}")

            y_cursor = y
            for para in (model.narrative.paragraphs or []):
                if _has_text(para):
                    logging.debug(f"Adding narrative para: '{para[:60]}...' at y={y_cursor}")
                    add_text(slide, x, y_cursor, w, 0.8, para, size=body_pt)
                    y_cursor += (0.28 if body_pt <= 16 else 0.32) + ST["narrative"]["para_gap_in"]

            if _has_list(model.narrative.bullets):
                logging.debug("Adding narrative bullets")
                b = add_bullets(slide, x, y_cursor, w, h - (y_cursor - y), model.narrative.bullets, size=bullet_pt)
                b.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT

        if has_kpi:
            logging.debug("Adding KPI tiles")
            kpi_top = L_top + (card_h + 0.25 if has_narr else 0.0)
            tile_w  = (L_w - 2 * kpi_gap) / 3.0
            tile_h  = kpi_h
            for i, k in enumerate((model.kpis or [])[:3]):
                if not _has_text(getattr(k, "headline", None)) and not _has_text(getattr(k, "caption", None)):
                    continue
                xk = L_left + i * (tile_w + kpi_gap)
                caption = wrap_text(k.caption, limit=ST["kpi"]["wrap_limit"])
                logging.debug(f"KPI {i}: x={xk}, top={kpi_top}, w={tile_w}, h={tile_h}, head='{k.headline}', cap='{caption}'")
                tile = add_kpi_tile(slide, xk, kpi_top, tile_w, tile_h,
                                    headline=k.headline, caption=caption, bg_hex=k.color_hex,
                                    headline_pt=ST["kpi"]["headline_pt"], caption_pt=ST["kpi"]["caption_pt"])
                tile.name = "KPI_TILE"

    # RIGHT — create only if blocks have content
    y_cursor = R_top

    has_steps = bool(model.steps and (_has_text(model.steps.header) or _has_list(model.steps.items)))
    logging.debug(f"RIGHT: has_steps={has_steps}")
    if has_steps:
        steps_h    = R_h * ST["steps"]["height_ratio"]
        logging.debug(f"Steps card height ratio -> {steps_h}")
        steps_card = add_card(slide, R_left, y_cursor, R_w, steps_h, radius=True, shadow=ST["shadow"]["enabled"])
        steps_card.name = "CARD_STEPS"
        if _has_text(model.steps.header):
            add_card_header(slide, steps_card, model.steps.header)
        if _has_list(model.steps.items):
            tb = add_bullets(slide, R_left+0.3, y_cursor+0.7, R_w-0.6, steps_h-0.9,
                             model.steps.items, size=ST["steps"]["items_pt"], numbered=True)
            tb.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        y_cursor += steps_h + ST["steps"]["gap_below_in"]
        logging.debug(f"RIGHT y_cursor -> {y_cursor}")

    has_icons = bool(model.icon_highlights and any(_has_text(i.caption) for i in model.icon_highlights))
    logging.debug(f"RIGHT: has_icons={has_icons}")
    if has_icons:
        gutter = ST["icons"]["gap_in"]
        tile_w = (R_w - 2*gutter) / 3.0
        tile_h = ST["icons"]["height_in"]
        logging.debug(f"Icon tiles: w={tile_w}, h={tile_h}")

        base_dir = os.path.abspath(os.path.join(here, ".."))  # v2/
        for i, icon in enumerate([i for i in model.icon_highlights if _has_text(i.caption)][:3]):
            x = R_left + i * (tile_w + gutter)
            logging.debug(f"Icon {i}: x={x}, y={y_cursor}, caption='{icon.caption}' icon_name='{icon.icon_name}'")
            card = add_card(slide, x, y_cursor, tile_w, tile_h, radius=True, shadow=True)
            card.name = "CARD_ICONS"

            img_rel = ICON_MAP.get(icon.icon_name or "", None)
            if img_rel:
                img_path = os.path.join(base_dir, img_rel)
                if os.path.exists(img_path):
                    logging.debug(f"Adding icon image: {img_path}")
                    slide.shapes.add_picture(img_path, Inches(x + tile_w/2 - ST["icons"]["img_h_in"]/2),
                                             Inches(y_cursor + 0.12), height=Inches(ST["icons"]["img_h_in"]))
                else:
                    logging.warning(f"Icon image not found: {img_path}")

            add_text(slide, x+0.2, y_cursor+0.8, tile_w-0.4, 0.6, icon.caption, size=ST["icons"]["caption_pt"])
        y_cursor += tile_h + ST["icons"]["gap_below_in"]
        logging.debug(f"RIGHT y_cursor -> {y_cursor}")

    has_outlook = bool(model.outlook and _has_text(model.outlook.text))
    logging.debug(f"RIGHT: has_outlook={has_outlook}")
    if has_outlook:
        outlook_h = max(ST["outlook"]["min_height_in"], (R_top + R_h) - y_cursor)
        logging.debug(f"Outlook height -> {outlook_h}")
        out_card  = add_card(slide, R_left, y_cursor, R_w, outlook_h, radius=True, shadow=ST["shadow"]["enabled"])
        out_card.name = "CARD_OUTLOOK"
        pad_l, pad_t, pad_r, pad_b = ST["outlook"]["padding_in"]
        add_text_padded(slide, R_left, y_cursor, R_w, outlook_h, model.outlook.text,
                        padding=(pad_l, pad_t, pad_r, pad_b), size=ST["outlook"]["body_pt"])

    # Footer — smaller + auto numbering
    if getattr(model, "footer", None):
        left_text  = model.footer.left_text
        prefix     = ST["footer"].get("prefix", "Slide ")
        right_text = f"{prefix}{slide_num}"
        logging.debug(f"Adding footer bar: left='{left_text}', right='{right_text}'")
        bar = add_footer_bar(
            slide,
            left=0.6,
            top=ST["page"]["height_in"] - (ST["footer"]["height_in"] + 0.35),
            width=ST["page"]["width_in"] - 1.2,
            height=ST["footer"]["height_in"],
            left_text=left_text,
            right_text=right_text,
            left_pt=ST["footer"]["left_pt"],
            right_pt=ST["footer"]["right_pt"],
        )
        bar.name = "FOOTER_BAR"

    logging.debug("Cleaning slide")
    cleanup_slide(slide)
    logging.info("=== render_slide end ===")
    return prs
