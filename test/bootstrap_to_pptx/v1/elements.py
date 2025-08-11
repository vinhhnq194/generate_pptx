# elements.py
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PRIMARY = RGBColor(13,110,253)
GRAY_900 = RGBColor(13,45,82)
GRAY_600 = RGBColor(108,117,125)
LIGHT_BG = RGBColor(248,249,250)
VERY_LIGHT = RGBColor(241,243,245)
WHITE = RGBColor(255,255,255)
DARK = RGBColor(33,37,41)

def add_title(slide, left, top, width, height, text):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GRAY_900
    return box

def add_subtitle(slide, left, top, width, height, text):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_600
    return box

def add_decor_diagonal(slide, left, top, width, height, angle=-10, rgba=(13,110,253, 15)):  # alpha ~ 6%
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.rotation = angle
    shape.fill.solid()
    r,g,b,a = rgba
    shape.fill.fore_color.rgb = RGBColor(r,g,b)
    shape.fill.transparency = 0.85
    shape.line.fill.background()
    return shape

def add_decor_circle(slide, left, top, diameter, rgba=(13,110,253, 15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    r,g,b,a = rgba
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r,g,b)
    shape.fill.transparency = 0.85
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, radius=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(222,226,230)  # subtle border
    shape.shadow.inherit = True
    return shape

def add_card_header(slide, card_shape, text, height_in=0.5, bg=PRIMARY, fg=WHITE):
    # Position header atop the card
    left = card_shape.left
    top = card_shape.top
    width = card_shape.width
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(height_in))
    header.fill.solid()
    header.fill.fore_color.rgb = bg
    header.line.fill.background()
    # text on header
    tb = slide.shapes.add_textbox(left+Inches(0.2), top+Inches(0.08), width-Inches(0.4), Inches(height_in-0.16))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = fg
    return header

def add_text(slide, left, top, width, height, text, size=14, color=RGBColor(52,58,64)):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    for i, para in enumerate(text.split("\n")):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = para
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return tb

def add_bullets(slide, left, top, width, height, items, size=14):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p._element.get_or_add_pPr().insert(0, tf._element._new_buChar())  # ensures bullets
    return tb

def add_kpi_tile(slide, left, top, width, height, headline, caption, bg=PRIMARY, fg=WHITE):
    tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    tile.fill.solid()
    tile.fill.fore_color.rgb = bg
    tile.line.fill.background()
    # headline
    tb = slide.shapes.add_textbox(tile.left+Inches(0.1), tile.top+Inches(0.1), tile.width-Inches(0.2), tile.height-Inches(0.2))
    tf = tb.text_frame
    tf.clear()
    p1 = tf.paragraphs[0]
    p1.text = headline
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = fg
    p2 = tf.add_paragraph()
    p2.text = caption
    p2.font.size = Pt(12)
    p2.font.color.rgb = fg
    return tile

def add_footer_bar(slide, left, top, width, height, left_text, right_text):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERY_LIGHT
    bar.line.color.rgb = RGBColor(222,226,230)
    tb_left = slide.shapes.add_textbox(bar.left+Inches(0.3), bar.top+Inches(0.05), Inches(width/2.0), Inches(height-0.1))
    tb_left.text_frame.text = left_text
    tb_left.text_frame.paragraphs[0].font.size = Pt(12)
    tb_left.text_frame.paragraphs[0].font.color.rgb = RGBColor(73,80,87)
    tb_right = slide.shapes.add_textbox(bar.left+bar.width-Inches(width/2.0)-Inches(0.3), bar.top+Inches(0.05),
                                        Inches(width/2.0), Inches(height-0.1))
    p = tb_right.text_frame.paragraphs[0]
    p.text = right_text
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(108,117,125)
    p.alignment = 2  # right
    return bar
