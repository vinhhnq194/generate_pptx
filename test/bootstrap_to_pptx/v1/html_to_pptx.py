# html_to_pptx.py
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches
from grid import Grid12
from elements import (
    add_title, add_subtitle, add_decor_diagonal, add_decor_circle,
    add_card, add_card_header, add_text, add_bullets,
    add_kpi_tile, add_footer_bar, PRIMARY, RGBColor
)

# --- 1) Load HTML ---
with open("/mnt/data/test.html", "r", encoding="utf-8") as f:
    soup = BeautifulSoup(f.read(), "html.parser")

prs = Presentation()
# 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

grid = Grid12(slide_width_in=13.333, margins_in=(0.6, 0.6, 0.6, 0.7), gutter_in=0.2)

# --- 2) Decorative shapes approximating CSS ---
# big diagonal band (oversized so rotation covers slide)
add_decor_diagonal(slide, left=-2.5, top=-2.0, width=16.0, height=10.0, angle=-10)
# circle glow bottom-right
add_decor_circle(slide, left=11.0, top=6.4, diameter=2.5)

# --- 3) Title & subtitle ---
title_el = soup.select_one("h2.fw-bold")
subtitle_el = soup.select_one("p.lead")
if title_el:
    add_title(slide, left=0.6, top=0.5, width=12.0, height=0.8, text=title_el.get_text(strip=True))
if subtitle_el:
    add_subtitle(slide, left=0.6, top=1.2, width=12.0, height=0.6, text=subtitle_el.get_text(strip=True))

# --- 4) Main two columns (.row > .col-6 / col-6) ---
main_row = soup.select_one(".row.g-3.flex-grow-1")
left_col = main_row.select_one(".col-6:nth-of-type(1)") if main_row else None
right_col = main_row.select_one(".col-6:nth-of-type(2)") if main_row else None

# Layout numbers
row_top = 1.8
row_height = 4.8

# Left column (narrative card + 3 KPI tiles)
if left_col:
    L_left, L_top, L_w, L_h = grid.rect_for(row_top_in=row_top, col_start=0, col_span=6, height_in=row_height)
    # Narrative card: take upper ~70% of left column
    card_h = L_h * 0.68
    card = add_card(slide, L_left, L_top, L_w, card_h)
    # Narrative text and bullets
    card_body = left_col.select_one(".card .card-body")
    if card_body:
        paras = [p.get_text(" ", strip=True) for p in card_body.select("p")]
        ul_items = [li.get_text(" ", strip=True) for li in card_body.select("ul li")]
        body_left = L_left + 0.3
        body_top  = L_top + 0.3
        body_w    = L_w   - 0.6
        body_h    = card_h - 0.6
        # Combine: first paragraph(s), then bullets
        y_cursor = body_top
        for para in paras:
            add_text(slide, body_left, y_cursor, body_w, 0.8, para, size=14)
            y_cursor += 0.7
        if ul_items:
            add_bullets(slide, body_left, y_cursor, body_w, 1.2, ul_items, size=14)

    # KPI row (3 tiles) below card
    kpi_top = L_top + card_h + 0.25
    gutter = 0.2
    tile_w = (L_w - 2 * gutter) / 3.0
    tile_h = L_h - (card_h + 0.35)

    # Extract KPI boxes
    stat_boxes = left_col.select(".row.g-2 .stat-box")
    kpis = []
    for sb in stat_boxes[:3]:
        headline = (sb.select_one(".fw-bold") or sb).get_text(" ", strip=True)
        caption = (sb.select_one("small") or sb).get_text(" ", strip=True)
        # background class decides color
        bg = sb.get("class", [])
        if "bg-secondary" in bg:
            color = RGBColor(108,117,125)
        elif "bg-dark" in bg:
            color = RGBColor(33,37,41)
        else:
            color = PRIMARY
        kpis.append((headline, caption, color))

    for i, (hd, cap, col) in enumerate(kpis):
        x = L_left + i * (tile_w + gutter)
        add_kpi_tile(slide, x, kpi_top, tile_w, tile_h, hd, cap, bg=col)

# Right column (steps card, 3 icon cards, outlook card)
if right_col:
    R_left, R_top, R_w, R_h = grid.rect_for(row_top_in=row_top, col_start=6, col_span=6, height_in=row_height)

    # Steps card (top)
    steps_card_h = R_h * 0.42
    steps_card = add_card(slide, R_left, R_top, R_w, steps_card_h)
    # header
    header_text = (right_col.select_one(".card.border-start .card-header") or {}).get_text(" ", strip=True) if right_col.select_one(".card.border-start") else "Three Critical Enablers"
    add_card_header(slide, steps_card, header_text)
    # ordered list
    ol_items = [li.get_text(" ", strip=True) for li in right_col.select(".card.border-start ol li")]
    if ol_items:
        add_bullets(slide, R_left+0.3, R_top+0.7, R_w-0.6, steps_card_h-0.9,
                    [f"{i+1}. {t}" for i,t in enumerate(ol_items)], size=14)

    # Icon highlights (3 cards), approximate with simple mini cards + captions
    icons_top = R_top + steps_card_h + 0.25
    icons_row = right_col.select(".row.g-2")[0] if right_col.select(".row.g-2") else None
    if icons_row:
        gutter = 0.2
        tile_w = (R_w - 2*gutter) / 3.0
        tile_h = 1.4
        captions = [cap.get_text(" ", strip=True) for cap in icons_row.select(".card p.small")]
        for i, caption in enumerate(captions[:3]):
            x = R_left + i * (tile_w + gutter)
            card = add_card(slide, x, icons_top, tile_w, tile_h, radius=True)
            add_text(slide, x+0.2, icons_top+0.4, tile_w-0.4, 0.6, caption, size=12, color=RGBColor(73,80,87))
            # NOTE: real Bootstrap icons aren't in HTML; replace with emoji or add_picture(...) if you have PNGs.

    # Outlook card (bottom, fills remaining space)
    outlook_top = icons_top + 1.6
    outlook_h = (R_top + R_h) - outlook_top
    outlook_card = add_card(slide, R_left, outlook_top, R_w, outlook_h)
    outlook_text = right_col.select(".card.flex-grow-1 .card-body p")
    if outlook_text:
        txt = " ".join(p.get_text(" ", strip=True) for p in outlook_text)
        add_text(slide, R_left+0.3, outlook_top+0.3, R_w-0.6, outlook_h-0.6, txt, size=14)

# --- 5) Footer bar ---
footer = soup.select_one(".footer-bar")
if footer:
    left_text = footer.select_one("div:nth-of-type(1)").get_text(" ", strip=True)
    right_text = footer.select_one("div:nth-of-type(2)").get_text(" ", strip=True)
    add_footer_bar(slide, left=0.0+0.6, top=7.5-0.9, width=13.333-1.2, height=0.5, left_text=left_text, right_text=right_text)

prs.save("consulting_slide.pptx")
print("Saved consulting_slide.pptx")
