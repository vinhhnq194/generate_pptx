from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from pptx.dml.color import RGBColor

GRAY_700 = RGBColor(0x51,0x51,0x51)
BORDER = RGBColor(0xDE,0xE2,0xE6)
VERY_LIGHT = RGBColor(0xF8,0xF9,0xFA)

def add_card(slide, left, top, width, height, radius=True, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid(); shp.fill.fore_color.rgb = RGBColor(0xFF,0xFF,0xFF)
    shp.line.color.rgb = BORDER
    shp.shadow.inherit = bool(shadow)
    if radius:
        try: shp.adjustments[0] = 0.16
        except: pass
    return shp

def add_title(slide, left, top, width, height, text):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(36); p.font.bold = True
    return tb

def add_subtitle(slide, left, top, width, height, text):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.paragraphs[0].text = text
    tf.paragraphs[0].font.size = Pt(17)
    return tb

def add_text(slide, left, top, width, height, text, size=14, color=GRAY_700):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    p = tf.paragraphs[0]; p.text = text or ""; p.font.size = Pt(size); p.font.color.rgb = color
    return tb

def add_bullets(slide, left, top, width, height, items, size=14, numbered=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE; tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, t in enumerate(items or []):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        prefix = f"{i+1}. " if numbered else "• "
        p.text = prefix + (t or "")
        p.font.size = Pt(size); p.font.color.rgb = GRAY_700
        p.space_after = Pt(4)
    return tb

def add_footer_bar(slide, left, top, width, height, left_text, right_text, left_pt=9, right_pt=9):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bar.fill.solid(); bar.fill.fore_color.rgb = VERY_LIGHT; bar.line.color.rgb = BORDER
    tb_left = slide.shapes.add_textbox(bar.left+Inches(0.28), bar.top+Inches(0.04), Inches(width/2), Inches(height-0.08))
    p = tb_left.text_frame.paragraphs[0]; p.text = left_text or ""; p.font.size = Pt(left_pt)
    tb_right = slide.shapes.add_textbox(bar.left+bar.width-Inches(width/2)-Inches(0.28), bar.top+Inches(0.04), Inches(width/2), Inches(height-0.08))
    r = tb_right.text_frame.paragraphs[0]; r.text = right_text or ""; r.font.size = Pt(right_pt); r.alignment = 2
    return bar

def add_kpi_tile(slide, left, top, width, height, headline, caption, bg_hex="#0d6efd", headline_pt=26, caption_pt=11):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    h = bg_hex.lstrip("#"); shp.fill.fore_color.rgb = RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    shp.line.fill.background()
    try: shp.adjustments[0] = 0.16
    except: pass
    add_text(slide, left+0.2, top+0.18, width-0.4, 0.5, headline, size=headline_pt, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(slide, left+0.2, top+0.70, width-0.4, 0.4, caption, size=caption_pt, color=RGBColor(0xF1,0xF3,0xF5))
    return shp
