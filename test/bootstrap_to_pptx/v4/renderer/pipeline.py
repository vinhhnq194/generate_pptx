import json
from pptx import Presentation
from pptx.util import Inches
from .grid import Grid12
from .layout_solver import solve_layout, Rect
from .element_registry import render_card, render_kpis, render_steps, render_icon_row, render_table, render_image, render_text
from utils.cleanup import cleanup_slide
from parsers.generic_bootstrap_to_ilt import parse_generic_bootstrap_to_ilt, ILTItem
from .elements import add_title, add_subtitle, add_footer_bar

def build_deck_from_html(html_path: str, styles_path: str, presets_path: str, template_path: str | None = None):
    ST = json.load(open(styles_path, "r", encoding="utf-8"))
    PRE = json.load(open(presets_path, "r", encoding="utf-8"))

    ilt = parse_generic_bootstrap_to_ilt(html_path)

    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width  = Inches(ST["page"]["width_in"])
    prs.slide_height = Inches(ST["page"]["height_in"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    grid = Grid12(
        slide_width_in=ST["page"]["width_in"],
        slide_height_in=ST["page"]["height_in"],
        margins_in=tuple(ST["page"]["margins_in"]),
        gutter_in=ST["page"]["gutter_in"]
    )

    # titles
    if ilt.title:    add_title(slide, 0.6, ST["title"]["top_in"], 12.0, 0.9, ilt.title)
    if ilt.subtitle: add_subtitle(slide, 0.6, ST["subtitle"]["top_in"], 12.0, 0.6, ilt.subtitle)

    placements = solve_layout(grid, ilt, row_top_in=ST["bands"]["row_top_in"], row_height_in=ST["bands"]["row_height_in"], row_gap_in=0.0)

    # render
    # (simple grouping: if adjacent items in same rect-row share kind 'kpi' or 'icon', you can batch render. Here: per item; KPI/icon use available width)
    for rect, it in placements:
        if it.kind == "card":   render_card(slide, rect, it, PRE, ST)
        elif it.kind == "kpi":  render_kpis(slide, Rect(rect.left, rect.top, rect.width, ST["kpi"]["height_in"]), [it], PRE, ST)
        elif it.kind == "steps":render_steps(slide, rect, it, PRE, ST)
        elif it.kind == "icon": render_icon_row(slide, Rect(rect.left, rect.top, rect.width, ST["icons"]["height_in"]), [it], PRE, ST)
        elif it.kind == "table":render_table(slide, rect, it, PRE)
        elif it.kind == "image":render_image(slide, rect, it, PRE)
        elif it.kind == "chart":render_text(slide, rect, ILTItem(kind="text", content={"text":"[Chart placeholder]"}), PRE)
        else:                   render_text(slide, rect, it, PRE)

    # footer
    if ilt.footer_left:
        add_footer_bar(slide,
            left=0.6,
            top=ST["page"]["height_in"] - (ST["footer"]["height_in"] + 0.40),
            width=ST["page"]["width_in"] - 1.2,
            height=ST["footer"]["height_in"],
            left_text=ilt.footer_left,
            right_text=f"{ST['footer'].get('prefix','Pg ')}1",
            left_pt=ST["footer"]["left_pt"],
            right_pt=ST["footer"]["right_pt"]
        )

    cleanup_slide(slide)
    return prs
