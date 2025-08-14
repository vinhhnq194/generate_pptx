# v4/renderer/element_registry.py
from typing import List
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR

from .elements import add_card, add_text, add_bullets, add_kpi_tile
from utils.bootstrap_mapping import apply_shape_appearance_from_bootstrap
from utils.text_fit import wrap_text


def _init_textframe(tb):
    """Standard text frame setup: wrap + shrink-to-fit + top anchor."""
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return tf


def render_card(slide, rect, item, presets: dict, styles: dict):
    """Render a narrative card (paragraphs + bullets) inside a rounded rectangle."""
    pad_l, pad_t, pad_r, pad_b = presets["card"]["padding_in"]
    x, y, w, h = rect.left, rect.top, rect.width, rect.height

    card = add_card(
        slide, x, y, w, h,
        radius=presets["card"]["rounded"],
        shadow=presets["card"]["shadow"]
    )
    apply_shape_appearance_from_bootstrap(card, item.classes)

    # One textbox for both paragraphs and bullets (more stable layout).
    tx, ty = x + pad_l, y + pad_t
    tw, th = w - (pad_l + pad_r), h - (pad_t + pad_b)
    tb = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
    tf = _init_textframe(tb)

    body_pt = presets["card"]["text"]["body_pt"]
    bullet_pt = presets["card"]["text"]["bullet_pt"]

    paras = item.content.get("paragraphs", []) or []
    bullets = item.content.get("bullets", []) or []

    # Paragraphs
    for i, para in enumerate(paras):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = para or ""
        p.font.size = Pt(body_pt)

    # Bullets
    for b in bullets:
        p = tf.add_paragraph()
        p.text = f"• {b or ''}"
        p.font.size = Pt(bullet_pt)
        p.space_after = Pt(3)


def render_kpis(slide, rect, items: List, presets: dict, styles: dict):
    """Render 1..N KPI tiles left-to-right within rect."""
    n = max(1, len(items))
    gutter = styles["kpi"]["gap_in"]
    tile_w = (rect.width - (n - 1) * gutter) / n
    tile_h = styles["kpi"]["height_in"]
    x, y = rect.left, rect.top

    for i, k in enumerate(items):
        headline = k.content.get("headline", "") or ""
        caption = wrap_text(k.content.get("caption", "") or "", limit=styles["kpi"]["wrap_limit"])
        add_kpi_tile(
            slide,
            x + i * (tile_w + gutter),
            y,
            tile_w,
            tile_h,
            headline=headline,
            caption=caption,
            bg_hex=presets["kpi"]["bg_hex"],
            headline_pt=presets["kpi"]["headline_pt"],
            caption_pt=presets["kpi"]["caption_pt"],
        )


def render_steps(slide, rect, item, presets: dict, styles: dict):
    """Render a steps card with a numbered list."""
    # Card container
    card = add_card(
        slide, rect.left, rect.top, rect.width, rect.height,
        radius=presets["steps"]["rounded"],
        shadow=presets["steps"]["shadow"]
    )

    # Optional header if provided
    header = (item.content.get("header") or "").strip()
    header_top = rect.top + 0.18
    bullets_top = rect.top + (0.18 if not header else 0.52)

    if header:
        tb = add_text(slide, rect.left + 0.3, header_top, rect.width - 0.6, 0.30, header, size=max(13, styles["steps"]["items_pt"]))
        # make header bold
        try:
            p = tb.text_frame.paragraphs[0]
            p.font.bold = True
        except Exception:
            pass

    # Numbered list
    add_bullets(
        slide,
        rect.left + 0.3,
        bullets_top,
        rect.width - 0.6,
        rect.height - (bullets_top - rect.top) - 0.22,
        item.content.get("items", []) or [],
        size=styles["steps"]["items_pt"],
        numbered=True
    )


def render_icon_row(slide, rect, items: List, presets: dict, styles: dict):
    """Render 1..N icon-caption tiles left-to-right."""
    n = max(1, len(items))
    gutter = styles["icons"]["gap_in"]
    tile_w = (rect.width - (n - 1) * gutter) / n
    tile_h = styles["icons"]["height_in"]
    x, y = rect.left, rect.top

    for i, it in enumerate(items):
        # Card per icon
        card = add_card(
            slide,
            x + i * (tile_w + gutter),
            y,
            tile_w,
            tile_h,
            radius=presets["icon"]["rounded"],
            shadow=presets["icon"]["shadow"]
        )
        # Caption
        add_text(
            slide,
            x + i * (tile_w + gutter) + 0.2,
            y + 0.78,
            tile_w - 0.4,
            0.6,
            it.content.get("caption", "") or "",
            size=presets["icon"]["caption_pt"]
        )


def render_table(slide, rect, item, presets: dict):
    """Render a simple table sized to rect."""
    rows = item.content.get("rows", []) or []
    if not rows:
        return
    table = slide.shapes.add_table(
        rows=len(rows),
        cols=len(rows[0]),
        left=Inches(rect.left),
        top=Inches(rect.top),
        width=Inches(rect.width),
        height=Inches(rect.height)
    ).table

    from pptx.dml.color import RGBColor

    for r, row in enumerate(rows):
        for c, txt in enumerate(row):
            cell = table.cell(r, c)
            cell.text = (txt or "")
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF3, 0xF5)


def render_image(slide, rect, item, presets: dict):
    """Render an image; maintain aspect by specifying width only."""
    src = item.content.get("src")
    if not src:
        return
    slide.shapes.add_picture(src, Inches(rect.left), Inches(rect.top), width=Inches(rect.width))


def render_text(slide, rect, item, presets: dict):
    """Render a plain text block."""
    add_text(
        slide,
        rect.left,
        rect.top,
        rect.width,
        rect.height,
        item.content.get("text", "") or "",
        size=presets["text"]["pt"]
    )
