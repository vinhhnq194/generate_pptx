import json, os
from pptx import Presentation
from pptx.util import Inches, Pt
from .grid import Grid12
from .layout_solver import Rect
from parsers.generic_bootstrap_to_ilt import parse_generic_bootstrap_to_ilt, ILTItem
from .element_registry import (
    render_card, render_kpis, render_steps, render_icon_row,
    render_table, render_image, render_text
)
from utils.merge import deep_update
from utils.clean_up import cleanup_slide

def layout_columns(grid: Grid12, ilt_rows, band_top_in: float, band_height_in: float, row_gap_in: float = 0.2):
    """Yield (row_index, [(Rect, [ILTItem or group])...]). Groups similar kinds."""
    out = []; cur_top = band_top_in
    for r in ilt_rows:
        placed = []
        cursor = 0
        for item in r.items:
            # group icon/kpi siblings into one row slot, others as single
            if item.kind in ("kpi","icon"):
                # gather siblings of same kind contiguous
                group = [item]
                # (a simple way; you could do smarter grouping outside)
                left, top, width, height = grid.rect_for(
                    row_top_in=cur_top - grid.m_top,
                    col_start=cursor + item.offset,
                    col_span=item.col_span,
                    height_in=band_height_in
                )
                placed.append((Rect(left, cur_top, width, band_height_in), group))
                cursor += item.offset + item.col_span
            else:
                left, top, width, height = grid.rect_for(
                    row_top_in=cur_top - grid.m_top,
                    col_start=cursor + item.offset,
                    col_span=item.col_span,
                    height_in=band_height_in
                )
                placed.append((Rect(left, cur_top, width, band_height_in), [item]))
                cursor += item.offset + item.col_span
        out.append((cur_top, placed))
        cur_top += band_height_in + row_gap_in
    return out

def build_deck_from_html(html_path: str, styles_path: str, presets_path: str, overrides_path: str = None, template_path: str = None):
    ST = json.load(open(styles_path,"r",encoding="utf-8"))
    PRE = json.load(open(presets_path,"r",encoding="utf-8"))
    if overrides_path and os.path.exists(overrides_path):
        deep_update(ST, json.load(open(overrides_path,"r",encoding="utf-8")))

    ilt = parse_generic_bootstrap_to_ilt(html_path)

    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width  = Inches(ST["page"]["width_in"])
    prs.slide_height = Inches(ST["page"]["height_in"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    grid = Grid12(
        slide_width_in = ST["page"]["width_in"],
        slide_height_in= ST["page"]["height_in"],
        margins_in     = tuple(ST["page"]["margins_in"]),
        gutter_in      = ST["page"]["gutter_in"]
    )

    # title/subtitle (optional)
    from .elements import add_title, add_subtitle, add_footer_bar
    if ilt.title:    add_title(slide, 0.6, ST["title"]["top_in"], 12.0, 0.9, ilt.title)
    if ilt.subtitle: add_subtitle(slide, 0.6, ST["subtitle"]["top_in"], 12.0, 0.6, ilt.subtitle)

    # compute rows
    band_top    = ST.get("bands", {}).get("row_top_in", 1.65)
    band_height = ST.get("bands", {}).get("row_height_in", 5.05)
    rows = layout_columns(grid, ilt.rows, band_top, band_height, row_gap_in=0.0)

    # render each placed slot
    for row_top, placed in rows:
        for rect, group in placed:
            kinds = [g.kind for g in group]
            # KPI / ICON are rendered as row groups
            if all(k == "kpi" for k in kinds):
                render_kpis(slide, rect, group, PRE, ST)
            elif all(k == "icon" for k in kinds):
                render_icon_row(slide, rect, group, PRE, ST)
            else:
                # single items
                it = group[0]
                # compute inner height if h_frac present
                inner_rect = Rect(rect.left, rect.top, rect.width,
                                  rect.height * (it.h_frac if it.h_frac else 1.0))
                if it.kind == "card":   render_card(slide, inner_rect, it, PRE, ST)
                elif it.kind == "steps":render_steps(slide, inner_rect, it, PRE, ST)
                elif it.kind == "table":render_table(slide, inner_rect, it, PRE)
                elif it.kind == "image":render_image(slide, inner_rect, it, PRE)
                elif it.kind == "chart":
                    # placeholder: you can implement your chart spec→pptx here
                    render_text(slide, inner_rect, ILTItem(kind="text", content={"text":"[Chart placeholder]"}), PRE)
                elif it.kind == "text": render_text(slide, inner_rect, it, PRE)
                else:
                    # fallback: render as text
                    render_text(slide, inner_rect, ILTItem(kind="text", content={"text": str(it.content)}), PRE)

    # footer (auto-number)
    if ilt.footer_left:
        add_footer_bar(slide, left=0.6,
                       top=ST["page"]["height_in"] - (ST["footer"]["height_in"] + 0.40),
                       width=ST["page"]["width_in"] - 1.2,
                       height=ST["footer"]["height_in"],
                       left_text=ilt.footer_left,
                       right_text=f"{ST['footer'].get('prefix','Pg ')}1",
                       left_pt=ST["footer"]["left_pt"],
                       right_pt=ST["footer"]["right_pt"])

    from utils.clean_up import cleanup_slide
    cleanup_slide(slide)
    return prs
