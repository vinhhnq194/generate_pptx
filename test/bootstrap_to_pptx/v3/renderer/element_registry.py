from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR
from .elements import add_card, add_card_header, add_text, add_bullets, add_kpi_tile
from utils.bootstrap_mapping import apply_shape_appearance_from_bootstrap
from utils.text_fit import wrap_text

def _tf_setup(tf):
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = MSO_ANCHOR.TOP

def render_card(slide, rect, item, presets, styles):
    pad = presets["card"]["padding_in"]
    x, y, w, h = rect.left, rect.top, rect.width, rect.height
    card = add_card(slide, x, y, w, h, radius=presets["card"]["rounded"], shadow=presets["card"]["shadow"])
    apply_shape_appearance_from_bootstrap(card, item.classes)

    tx, ty = x+pad[0], y+pad[1]
    tw, th = w-(pad[0]+pad[2]), h-(pad[1]+pad[3])

    tb = slide.shapes.add_textbox(Inches(tx), Inches(ty), Inches(tw), Inches(th))
    tf = tb.text_frame; _tf_setup(tf)
    # paragraphs
    for i, para in enumerate(item.content.get("paragraphs", [])):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = para; p.font.size = Pt(presets["card"]["text"]["body_pt"])
    # bullets
    for b in item.content.get("bullets", []):
        p = tf.add_paragraph(); p.text = f"• {b}"
        p.font.size = Pt(presets["card"]["text"]["bullet_pt"])

def render_kpis(slide, rect, items, presets, styles):
    gutter = styles["kpi"]["gap_in"]
    w = rect.width; tile_w = (w - (len(items)-1)*gutter) / max(1,len(items))
    tile_h = styles["kpi"]["height_in"]; x = rect.left; y = rect.top
    for i, k in enumerate(items):
        cap = wrap_text(k.content.get("caption",""), limit=styles["kpi"]["wrap_limit"])
        add_kpi_tile(slide, x + i*(tile_w + gutter), y, tile_w, tile_h,
                     headline=k.content.get("headline",""), caption=cap,
                     bg_hex=presets["kpi"]["bg_hex"],
                     headline_pt=presets["kpi"]["headline_pt"],
                     caption_pt=presets["kpi"]["caption_pt"])

def render_steps(slide, rect, item, presets, styles):
    card = add_card(slide, rect.left, rect.top, rect.width, rect.height,
                    radius=presets["steps"]["rounded"], shadow=presets["steps"]["shadow"])
    header = item.content.get("header") or "Steps"
    add_card_header(slide, card, header)
    tb = add_bullets(slide, rect.left+0.3, rect.top+0.7, rect.width-0.6, rect.height-0.9,
                     item.content.get("items", []),
                     size=presets["steps"]["item_pt"], numbered=True)

def render_icon_row(slide, rect, items, presets, styles):
    gutter = styles["icons"]["gap_in"]
    tile_w = (rect.width - (len(items)-1)*gutter)/max(1,len(items))
    tile_h = styles["icons"]["height_in"]; x=rect.left; y=rect.top
    for i, it in enumerate(items):
        card = add_card(slide, x + i*(tile_w+gutter), y, tile_w, tile_h,
                        radius=presets["icon"]["rounded"], shadow=presets["icon"]["shadow"])
        add_text(slide, x + i*(tile_w+gutter) + 0.2, y + 0.8, tile_w-0.4, 0.6,
                 it.content.get("caption",""), size=presets["icon"]["caption_pt"])

def render_table(slide, rect, item, presets):
    rows = item.content.get("rows", [])
    if not rows: return
    # simple table: let pptx compute cell sizes; we set overall bounds
    table = slide.shapes.add_table(rows=len(rows), cols=len(rows[0]),
                                   left=Inches(rect.left), top=Inches(rect.top),
                                   width=Inches(rect.width), height=Inches(rect.height)).table
    for r, row in enumerate(rows):
        for c, txt in enumerate(row):
            cell = table.cell(r, c)
            cell.text = txt or ""
            if r == 0:  # header
                cell.fill.solid(); 
                from pptx.dml.color import RGBColor
                cell.fill.fore_color.rgb = RGBColor(0xF1,0xF3,0xF5)

def render_image(slide, rect, item, presets):
    src = item.content.get("src")
    if not src: return
    slide.shapes.add_picture(src, Inches(rect.left), Inches(rect.top),
                             width=Inches(rect.width))  # keep aspect

def render_text(slide, rect, item, presets):
    add_text(slide, rect.left, rect.top, rect.width, rect.height,
             item.content.get("text",""), size=presets["text"]["pt"])
