import logging
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from utils.colors import hex_to_rgb_color

# Configure logging
logging.basicConfig(level=logging.DEBUG, format="%(levelname)s: %(message)s")

PRIMARY   = RGBColor(13,110,253)
GRAY_900  = RGBColor(13,45,82)
GRAY_700  = RGBColor(73,80,87)
GRAY_600  = RGBColor(108,117,125)
VERY_LIGHT= RGBColor(241,243,245)
BORDER    = RGBColor(222,226,230)
WHITE     = RGBColor(255,255,255)
DARK      = RGBColor(33,37,41)

def add_title(slide, left, top, width, height, text):
    logging.debug(f"Adding title at ({left},{top}), size=({width}x{height}), text='{text}'")
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = box.text_frame.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GRAY_900
    return box

def add_subtitle(slide, left, top, width, height, text):
    logging.debug(f"Adding subtitle at ({left},{top}), size=({width}x{height}), text='{text}'")
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    p = box.text_frame.paragraphs[0]
    p.text = text or ""
    p.font.size = Pt(18)
    p.font.color.rgb = GRAY_600
    return box

def add_decor_diagonal(slide, left, top, width, height, angle=-10, rgba=(13,110,253, 0.85)):
    logging.debug(f"Adding diagonal decor rect at ({left},{top}), size=({width}x{height}), angle={angle}, rgba={rgba}")
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.rotation = angle
    shape.fill.solid()
    r,g,b,alpha = rgba
    shape.fill.fore_color.rgb = RGBColor(r,g,b)
    shape.fill.transparency = alpha
    shape.line.fill.background()
    return shape

def add_decor_circle(slide, left, top, diameter, rgba=(13,110,253, 0.85)):
    logging.debug(f"Adding decor circle at ({left},{top}), diameter={diameter}, rgba={rgba}")
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    r,g,b,alpha = rgba
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r,g,b)
    shape.fill.transparency = alpha
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, radius=True, shadow=True):
    logging.debug(f"Adding card at ({left},{top}), size=({width}x{height}), radius={radius}, shadow={shadow}")
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER
    shape.shadow.inherit = shadow
    if radius and hasattr(shape, "adjustments"):
        try:
            shape.adjustments[0] = 0.15
            logging.debug("Applied softer corner radius to card.")
        except Exception as e:
            logging.warning(f"Could not set radius adjustment: {e}")
    return shape

def add_card_header(slide, card_shape, text, height_in=0.5, bg=PRIMARY, fg=WHITE):
    logging.debug(f"Adding card header text='{text}', height={height_in}, bg={bg}, fg={fg}")
    left = card_shape.left
    top = card_shape.top
    width = card_shape.width
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(height_in))
    header.fill.solid()
    header.fill.fore_color.rgb = bg
    header.line.fill.background()

    tb = slide.shapes.add_textbox(left+Inches(0.2), top+Inches(0.08), width-Inches(0.4), Inches(height_in-0.16))
    p = tb.text_frame.paragraphs[0]
    p.text = text or ""
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = fg
    return header

# def add_text(slide, left, top, width, height, text, size=14, color=GRAY_700):
#     logging.debug(f"Adding text at ({left},{top}), size=({width}x{height}), size_pt={size}, color={color}, text='{text}'")
#     tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
#     tf = tb.text_frame
#     tf.clear()
#     tf.word_wrap = True
#     for i, para in enumerate((text or "").split("\n")):
#         p = tf.add_paragraph() if i else tf.paragraphs[0]
#         p.text = para
#         p.font.size = Pt(size)
#         p.font.color.rgb = color
#     return tb

def add_text_padded(slide, outer_left, outer_top, outer_width, outer_height, text,
                    padding=(0.3,0.3,0.3,0.3), size=14, color=GRAY_700):
    logging.debug(f"Adding padded text with padding={padding}")
    pl, pt, pr, pb = padding
    return add_text(slide, outer_left+pl, outer_top+pt, outer_width-(pl+pr), outer_height-(pt+pb),
                    text, size=size, color=color)

# def add_bullets(slide, left, top, width, height, items, size=14, numbered=False):
#     logging.debug(f"Adding bullets at ({left},{top}), size=({width}x{height}), count={len(items)}, numbered={numbered}")
#     tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
#     tf = tb.text_frame
#     tf.clear()
#     tf.word_wrap = True

#     if not items:
#         logging.debug("No bullet items provided.")
#         return tb

#     for i, item in enumerate(items):
#         p = tf.add_paragraph() if i else tf.paragraphs[0]
#         prefix = f"{i+1}. " if numbered else "• "
#         p.text = prefix + (item or "")
#         p.level = 0
#         p.space_after = Pt(6)
#         p.font.size = Pt(size)
#         p.font.color.rgb = GRAY_700
#         logging.debug(f"Added bullet: '{p.text}'")

#     return tb

def add_kpi_tile(slide, left, top, width, height, headline, caption, bg_hex="#0d6efd",
                 headline_pt=28, caption_pt=12):
    logging.debug(f"Adding KPI tile at ({left},{top}), size=({width}x{height}), headline='{headline}', caption='{caption}', bg={bg_hex}")
    tile = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    tile.fill.solid()
    tile.fill.fore_color.rgb = hex_to_rgb_color(bg_hex)
    tile.line.fill.background()

    tb = slide.shapes.add_textbox(tile.left+Inches(0.12), tile.top+Inches(0.12),
                                  tile.width-Inches(0.24), tile.height-Inches(0.24))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.text = headline or ""
    p1.font.size = Pt(headline_pt)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    p2.text = caption or ""
    p2.font.size = Pt(caption_pt)
    p2.font.color.rgb = WHITE
    return tile

def add_footer_bar(slide, left, top, width, height, left_text, right_text,
                   left_pt=10, right_pt=10):
    logging.debug(f"Adding footer bar at ({left},{top}), size=({width}x{height}), left_text='{left_text}', right_text='{right_text}'")
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    bar.name = "FOOTER_BAR"
    bar.fill.solid()
    bar.fill.fore_color.rgb = VERY_LIGHT
    bar.line.color.rgb = BORDER

    tb_left = slide.shapes.add_textbox(bar.left+Inches(0.28), bar.top+Inches(0.04),
                                       Inches(width/2.0), Inches(height-0.08))
    l = tb_left.text_frame.paragraphs[0]
    l.text = left_text or ""
    l.font.size = Pt(left_pt)
    l.font.color.rgb = GRAY_700

    tb_right = slide.shapes.add_textbox(bar.left+bar.width-Inches(width/2.0)-Inches(0.28),
                                        bar.top+Inches(0.04),
                                        Inches(width/2.0), Inches(height-0.08))
    r = tb_right.text_frame.paragraphs[0]
    r.text = right_text or ""
    r.font.size = Pt(right_pt)
    r.font.color.rgb = GRAY_600
    r.alignment = 2  # right
    return bar

from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR

def add_text(sh, left, top, width, height, text, size=14, color=GRAY_700):
    tb = sh.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink text, don't grow shape
    tf.vertical_anchor = MSO_ANCHOR.TOP
    # tighter inner padding
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    for i, para in enumerate((text or "").split("\n")):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        p.text = para
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return tb

def add_bullets(slide, left, top, width, height, items, size=14, numbered=False):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # shrink to avoid overflow
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)

    if not items:
        return tb

    for i, item in enumerate(items):
        p = tf.add_paragraph() if i else tf.paragraphs[0]
        prefix = f"{i+1}. " if numbered else "• "
        p.text = prefix + (item or "")
        p.level = 0
        p.space_after = Pt(4)
        p.font.size = Pt(size)
        p.font.color.rgb = GRAY_700
    return tb
