# -*- coding: utf-8 -*-
"""
Bootstrap → python-pptx conversion helpers + mapping registry
Covers each line in the user's python_pptx_mutable_properties table.
"""

from dataclasses import dataclass
from typing import Any, Dict, List, Optional, Tuple, Callable, Union
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Inches, Cm, Pt

# -----------------------------
# Basic palettes / helpers
# -----------------------------

BOOTSTRAP_COLORS = {
    "primary":  RGBColor(0x0D, 0x6E, 0xFD),
    "secondary":RGBColor(0x6C, 0x75, 0x7D),
    "success":  RGBColor(0x19, 0x87, 0x54),
    "danger":   RGBColor(0xDC, 0x35, 0x45),
    "warning":  RGBColor(0xFF, 0xC1, 0x07),
    "info":     RGBColor(0x0D, 0xCA, 0xF0),
    "light":    RGBColor(0xF8, 0xF9, 0xFA),
    "dark":     RGBColor(0x21, 0x25, 0x29),
}

BORDER_GRAY = RGBColor(0xDE, 0xE2, 0xE6)

REM_MAP = {0:0.0, 1:0.25, 2:0.5, 3:1.0, 4:1.5, 5:3.0}  # rem units from Bootstrap spacers
def apply_shape_appearance_from_bootstrap(shape, class_list: List[str]):
    """
    Apply only visual styling from Bootstrap-like classes (no geometry).
    Safe for ILT pipeline where positions/sizes are already solved.
    """
    p = parse_bootstrap_utils(class_list)

    # Fill color
    if p.bg:
        convert_shape_fill_color(shape, p.bg)

    # Border
    if p.border_zero:
        try:
            shape.line.fill.background()
        except Exception:
            pass
    else:
        convert_shape_line_color_default(shape)
        convert_shape_line_width(shape, p.border)

    # Rounding
    if p.rounded:
        convert_shape_corner_radius(shape, 0.2)

def rem_to_inches(rem: float, px_per_inch: int = 96, base_font_px: int = 16) -> float:
    """Convert rem to inches, assuming 1rem = base_font_px pixels."""
    px = rem * base_font_px
    return px / px_per_inch

def percent_of(value: float, pct: float) -> float:
    return value * pct

# -----------------------------
# Parsing Bootstrap utilities
# -----------------------------

@dataclass
class ParsedUtils:
    col: Optional[int] = None           # col-1..12
    offset: int = 0                     # offset-0..11
    mt: Optional[int] = None            # mt-0..5
    bg: Optional[str] = None            # bg-primary, etc.
    text_color: Optional[str] = None    # text-primary, etc.
    border: Optional[int] = None        # border-1..5; None means default
    border_zero: bool = False           # border-0
    rounded: bool = False               # rounded
    fw_bold: bool = False               # fw-bold
    fst_italic: bool = False            # fst-italic
    text_uppercase: bool = False        # text-uppercase
    h_frac: Optional[float] = None      # h-25/50/75/100

def parse_bootstrap_utils(classes: List[str]) -> ParsedUtils:
    p = ParsedUtils()
    for c in classes or []:
        if c.startswith("col-"):
            try: p.col = int(c.split("-")[1])
            except: pass
        elif c.startswith("offset-"):
            try: p.offset = int(c.split("-")[1])
            except: pass
        elif c.startswith("mt-"):
            try: p.mt = int(c.split("-")[1])
            except: pass
        elif c.startswith("bg-"):
            key = c.replace("bg-","").strip()
            if key in BOOTSTRAP_COLORS: p.bg = key
        elif c.startswith("text-"):
            key = c.replace("text-","").strip()
            if key in BOOTSTRAP_COLORS: p.text_color = key
        elif c == "border-0":
            p.border_zero = True
        elif c.startswith("border-"):
            try: p.border = int(c.split("-")[1])
            except: pass
        elif c == "rounded":
            p.rounded = True
        elif c == "fw-bold":
            p.fw_bold = True
        elif c == "fst-italic":
            p.fst_italic = True
        elif c == "text-uppercase":
            p.text_uppercase = True
        elif c.startswith("h-"):
            frac = c.split("-")[1]
            frac_map = {"25":0.25,"50":0.5,"75":0.75,"100":1.0}
            if frac in frac_map:
                p.h_frac = float(frac_map[frac])
    return p

# -----------------------------
# Grid context (for width/left)
# -----------------------------

@dataclass
class GridContext:
    slide_width_in: float
    slide_height_in: float
    margin_left_in: float
    margin_right_in: float
    gutter_in: float

    @property
    def content_width_in(self) -> float:
        return self.slide_width_in - self.margin_left_in - self.margin_right_in

    @property
    def col_width_in(self) -> float:
        # 12 columns = 12 widths + 11 gutters
        return (self.content_width_in - 11*self.gutter_in) / 12.0

    def left_for(self, offset_cols: int) -> float:
        return self.margin_left_in + offset_cols * (self.col_width_in + self.gutter_in)

    def width_for(self, span_cols: int) -> float:
        if span_cols <= 0: return 0.0
        return span_cols * self.col_width_in + (span_cols - 1) * self.gutter_in

# -----------------------------
# Mapping Registry dataclass
# -----------------------------

@dataclass
class MappingRule:
    name: str
    variable: str
    category: str
    prop: str
    pptx_attr: str
    value_type_notes: str
    bootstrap_class: str
    apply: Callable[..., None]  # function implementing conversion

# -----------------------------
# Converters for each table row
# -----------------------------
# Slide background

def convert_slide_bg_color(slide, rgb: RGBColor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb

def convert_slide_bg_type(slide, kind: str):
    # 'solid' handled; 'gradient/patterned' best set on template or via drawn shapes.
    if kind == "solid":
        slide.background.fill.solid()
    else:
        # no-op; recommend template or overlay shape strategy
        pass

# Shape geometry

def convert_shape_left_from_offset(shape, grid: GridContext, offset_n: int):
    left_in = grid.left_for(offset_n)
    shape.left = Inches(left_in)

def convert_shape_top_from_mt(shape, mt_n: int):
    # Convert Bootstrap mt-n to inches via rem mapping
    rem = REM_MAP.get(mt_n, 0.0)
    top_in = rem_to_inches(rem)
    shape.top = Inches(top_in)

def convert_shape_width_from_col(shape, grid: GridContext, col_n: int):
    w_in = grid.width_for(col_n)
    shape.width = Inches(w_in)

def convert_shape_height_from_hfrac(shape, band_height_in: float, h_frac: float):
    shape.height = Inches(band_height_in * h_frac)

# Shape appearance

def convert_shape_fill_color(shape, bootstrap_bg: Optional[str]):
    if bootstrap_bg is None: return
    shape.fill.solid()
    shape.fill.fore_color.rgb = BOOTSTRAP_COLORS[bootstrap_bg]

def convert_shape_fill_type_solid(shape):
    shape.fill.solid()

def convert_shape_line_color_default(shape):
    shape.line.color.rgb = BORDER_GRAY

def convert_shape_line_width(shape, border_n: Optional[int]):
    if border_n is None: return
    width_map = {1:0.75, 2:1.0, 3:1.5, 4:2.0, 5:3.0}
    pt = width_map.get(border_n, 0.75)
    shape.line.width = Pt(pt)

def convert_shape_line_dash(shape, style: str):
    # style options: "solid", "dash", "dot", "dash_dot"
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    m = {
        "solid": MSO_LINE_DASH_STYLE.SOLID,
        "dash":  MSO_LINE_DASH_STYLE.DASH,
        "dot":   MSO_LINE_DASH_STYLE.DOT,
        "dash_dot": MSO_LINE_DASH_STYLE.DASH_DOT,
    }
    if style in m:
        shape.line.dash_style = m[style]

def convert_shape_shadow_enable(shape, on: bool):
    shape.shadow.inherit = bool(on)

def convert_shape_shadow_blur(shape, pt: float):
    shape.shadow.blur_radius = Pt(pt)

def convert_shape_shadow_distance(shape, pt: float):
    shape.shadow.distance = Pt(pt)

def convert_shape_shadow_direction(shape, degrees: float):
    shape.shadow.direction = float(degrees)

def convert_shape_corner_radius(shape, fraction_0_1: float = 0.2):
    # Works for rounded rectangle shapes
    try:
        shape.adjustments[0] = fraction_0_1
    except Exception:
        pass

# Text (shape.text_frame / paragraphs / runs)

def convert_text_content(shape, text: str):
    shape.text = text or ""

def convert_paragraph_alignment(paragraph, align_cls: str):
    m = {"start": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "end": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    paragraph.alignment = m.get(align_cls, PP_ALIGN.LEFT)

def convert_paragraph_space_before(paragraph, pt: float):
    paragraph.space_before = Pt(pt)

def convert_paragraph_space_after(paragraph, pt: float):
    paragraph.space_after = Pt(pt)

def convert_paragraph_line_spacing(paragraph, value: Union[float, int]):
    # float (multiplier) or Pt
    if isinstance(value, (int, float)) and value <= 3.0:
        paragraph.line_spacing = float(value)
    else:
        paragraph.line_spacing = Pt(float(value))

def convert_run_text(run, text: str):
    run.text = text or ""

def convert_run_bold(run, is_bold: bool):
    run.font.bold = bool(is_bold)

def convert_run_italic(run, is_italic: bool):
    run.font.italic = bool(is_italic)

def convert_run_size(run, pt: float):
    run.font.size = Pt(pt)

def convert_run_name(run, name: str):
    run.font.name = name

def convert_run_color(run, bootstrap_text: Optional[str]):
    if bootstrap_text is None: return
    run.font.color.rgb = BOOTSTRAP_COLORS[bootstrap_text]

def convert_run_underline(run, on: bool):
    run.font.underline = bool(on)

def convert_run_all_caps(run, on: bool):
    run.font.all_caps = bool(on)

# Table

def convert_cell_text(cell, text: str):
    cell.text = text or ""

def convert_cell_fill_color(cell, rgb: RGBColor):
    cell.fill.solid()
    cell.fill.fore_color.rgb = rgb

def convert_cell_merge(cell, other_cell):
    cell.merge(other_cell)

def convert_cell_text_alignment(cell, align: str):
    m = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    cell.text_frame.paragraphs[0].alignment = m.get(align, PP_ALIGN.LEFT)

# Image

def convert_add_picture(shapes, path: str, left_in: float, top_in: float, width_in: Optional[float]=None, height_in: Optional[float]=None):
    left, top = Inches(left_in), Inches(top_in)
    if width_in:
        return shapes.add_picture(path, left, top, width=Inches(width_in))
    if height_in:
        return shapes.add_picture(path, left, top, height=Inches(height_in))
    return shapes.add_picture(path, left, top)

# Chart

def convert_chart_legend(chart, has_legend: bool):
    chart.has_legend = bool(has_legend)

def convert_chart_legend_position(chart, pos: str):
    from pptx.enum.chart import XL_LEGEND_POSITION
    m = {
        "top": XL_LEGEND_POSITION.TOP,
        "bottom": XL_LEGEND_POSITION.BOTTOM,
        "left": XL_LEGEND_POSITION.LEFT,
        "right": XL_LEGEND_POSITION.RIGHT,
        "corner": XL_LEGEND_POSITION.CORNER,
    }
    if pos in m: chart.legend.position = m[pos]

def convert_chart_value_axis_max(chart, value: float):
    chart.value_axis.maximum_scale = float(value)

def convert_chart_value_axis_min(chart, value: float):
    chart.value_axis.minimum_scale = float(value)

def convert_series_name(series, name: str):
    series.name = name or ""

def convert_series_values(series, values: List[float]):
    series.values = values

def convert_series_x_values(series, x_values: List[Any]):
    series.x_values = x_values

def convert_series_fill_color(series, rgb: RGBColor):
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = rgb

# Notes

def convert_notes_text(slide, text: str):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text or ""

# Other

def convert_shape_hyperlink(shape, url: str):
    shape.click_action.hyperlink.address = url

def convert_shape_rotation(shape, degrees: float):
    shape.rotation = float(degrees)

# -----------------------------
# Mapping registry (for DB)
# Each entry matches a table row with a callable that performs the conversion.
# -----------------------------

MAPPING_REGISTRY: List[MappingRule] = [
    # Slide
    MappingRule("Slide BG Color", "slide_bg_color", "Slide", "Background Fill Color",
                "slide.background.fill.fore_color.rgb", "RGB tuple",
                "#N/A", lambda slide, rgb: convert_slide_bg_color(slide, rgb)),

    MappingRule("Slide BG Type", "slide_bg_type", "Slide", "Background Fill Type",
                "slide.background.fill.solid()", "solid, gradient, patterned",
                "#N/A", lambda slide, kind: convert_slide_bg_type(slide, kind)),

    # Shape geometry
    MappingRule("Shape Left", "shape_left", "Shape", "Left Position",
                "shape.left", "Length", "offset-1..offset-11",
                lambda shape, grid, offset_n: convert_shape_left_from_offset(shape, grid, offset_n)),

    MappingRule("Shape Top", "shape_top", "Shape", "Top Position",
                "shape.top", "Length", "mt-0..mt-5",
                lambda shape, mt_n: convert_shape_top_from_mt(shape, mt_n)),

    MappingRule("Shape Width", "shape_width", "Shape", "Width",
                "shape.width", "Length", "col-1..col-12",
                lambda shape, grid, col_n: convert_shape_width_from_col(shape, grid, col_n)),

    MappingRule("Shape Height", "shape_height", "Shape", "Height",
                "shape.height", "Length", "h-25",
                lambda shape, band_h_in, frac: convert_shape_height_from_hfrac(shape, band_h_in, frac)),

    # Shape appearance
    MappingRule("Shape Fill Color", "shape_fill_color", "Shape", "Fill Color",
                "shape.fill.fore_color.rgb", "RGB tuple", "bg-*",
                lambda shape, bg_key: convert_shape_fill_color(shape, bg_key)),

    MappingRule("Shape Fill Type", "shape_fill_type", "Shape", "Fill Type",
                "shape.fill.solid()", "solid/gradient/patterned", "#N/A",
                lambda shape, *_: convert_shape_fill_type_solid(shape)),

    MappingRule("Shape Line Color", "shape_line_color", "Shape", "Line Color",
                "shape.line.color.rgb", "RGB tuple", "border",
                lambda shape: convert_shape_line_color_default(shape)),

    MappingRule("Shape Line Width", "shape_line_width", "Shape", "Line Width",
                "shape.line.width", "Pt", "border-1..5",
                lambda shape, n: convert_shape_line_width(shape, n)),

    MappingRule("Shape Line Dash", "shape_line_dash", "Shape", "Line Dash Style",
                "shape.line.dash_style", "dash/solid/etc", "#N/A",
                lambda shape, style: convert_shape_line_dash(shape, style)),

    MappingRule("Shape Shadow Enable", "shape_shadow_on", "Shape", "Shadow Enable",
                "shape.shadow.inherit", "True/False", "#N/A",
                lambda shape, on: convert_shape_shadow_enable(shape, on)),

    MappingRule("Shape Shadow Blur", "shape_shadow_blur", "Shape", "Shadow Blur Radius",
                "shape.shadow.blur_radius", "Pt", "#N/A",
                lambda shape, pt: convert_shape_shadow_blur(shape, pt)),

    MappingRule("Shape Shadow Distance", "shape_shadow_distance", "Shape", "Shadow Distance",
                "shape.shadow.distance", "Pt", "#N/A",
                lambda shape, pt: convert_shape_shadow_distance(shape, pt)),

    MappingRule("Shape Shadow Direction", "shape_shadow_direction", "Shape", "Shadow Direction",
                "shape.shadow.direction", "Degrees", "#N/A",
                lambda shape, deg: convert_shape_shadow_direction(shape, deg)),

    MappingRule("Shape Corner Radius", "shape_corner_radius", "Shape", "Corner Radius",
                "shape.adjustments[0]", "Float 0-1", "rounded",
                lambda shape, frac=0.2: convert_shape_corner_radius(shape, frac)),

    # Text (shape level content)
    MappingRule("Text Content", "text_content", "Text", "Text Content",
                "shape.text", "String", "#N/A",
                lambda shape, text: convert_text_content(shape, text)),

    # Paragraph
    MappingRule("Paragraph Align", "para_align", "Text", "Paragraph Alignment",
                "paragraph.alignment", "LEFT/CENTER/RIGHT/JUSTIFY", "text-start",
                lambda paragraph, align: convert_paragraph_alignment(paragraph, align)),

    MappingRule("Paragraph Space Before", "para_space_before", "Text", "Paragraph Space Before",
                "paragraph.space_before", "Pt", "#N/A",
                lambda paragraph, pt: convert_paragraph_space_before(paragraph, pt)),

    MappingRule("Paragraph Space After", "para_space_after", "Text", "Paragraph Space After",
                "paragraph.space_after", "Pt", "#N/A",
                lambda paragraph, pt: convert_paragraph_space_after(paragraph, pt)),

    MappingRule("Paragraph Line Spacing", "para_line_spacing", "Text", "Paragraph Line Spacing",
                "paragraph.line_spacing", "Float or Pt", "#N/A",
                lambda paragraph, val: convert_paragraph_line_spacing(paragraph, val)),

    # Runs
    MappingRule("Run Text", "run_text", "Text", "Run Text",
                "run.text", "String", "#N/A",
                lambda run, text: convert_run_text(run, text)),

    MappingRule("Font Bold", "font_bold", "Text", "Font Bold",
                "run.font.bold", "True/False", "fw-bold",
                lambda run, is_bold: convert_run_bold(run, is_bold)),

    MappingRule("Font Italic", "font_italic", "Text", "Font Italic",
                "run.font.italic", "True/False", "fst-italic",
                lambda run, is_italic: convert_run_italic(run, is_italic)),

    MappingRule("Font Size", "font_size", "Text", "Font Size",
                "run.font.size", "Pt", "#N/A",
                lambda run, pt: convert_run_size(run, pt)),

    MappingRule("Font Name", "font_name", "Text", "Font Name",
                "run.font.name", "String", "#N/A",
                lambda run, name: convert_run_name(run, name)),

    MappingRule("Font Color", "font_color", "Text", "Font Color",
                "run.font.color.rgb", "RGB tuple", "text-primary",
                lambda run, text_key: convert_run_color(run, text_key)),

    MappingRule("Font Underline", "font_underline", "Text", "Font Underline",
                "run.font.underline", "True/False", "#N/A",
                lambda run, on: convert_run_underline(run, on)),

    MappingRule("Font All Caps", "font_all_caps", "Text", "Font All Caps",
                "run.font.all_caps", "True/False", "text-uppercase",
                lambda run, on: convert_run_all_caps(run, on)),

    # Table
    MappingRule("Cell Text", "cell_text", "Table", "Cell Text",
                "cell.text", "String", "#N/A",
                lambda cell, text: convert_cell_text(cell, text)),

    MappingRule("Cell Fill Color", "cell_fill_color", "Table", "Cell Fill Color",
                "cell.fill.fore_color.rgb", "RGB tuple", "#N/A",
                lambda cell, rgb: convert_cell_fill_color(cell, rgb)),

    MappingRule("Cell Merge", "cell_merge", "Table", "Cell Merge",
                "cell.merge()", "Cell object", "#N/A",
                lambda cell, other: convert_cell_merge(cell, other)),

    MappingRule("Cell Align", "cell_align", "Table", "Cell Text Alignment",
                "cell.text_frame.paragraphs[0].alignment", "LEFT/CENTER/RIGHT", "#N/A",
                lambda cell, align: convert_cell_text_alignment(cell, align)),

    # Image
    MappingRule("Add Picture", "add_picture", "Image", "Add Picture",
                "shapes.add_picture()", "path + pos/size", "#N/A",
                lambda shapes, path, left_in, top_in, width_in=None, height_in=None:
                    convert_add_picture(shapes, path, left_in, top_in, width_in, height_in)),

    # Chart
    MappingRule("Chart Legend", "chart_has_legend", "Chart", "Has Legend",
                "chart.has_legend", "True/False", "#N/A",
                lambda chart, on: convert_chart_legend(chart, on)),

    MappingRule("Legend Position", "chart_legend_position", "Chart", "Legend Position",
                "chart.legend.position", "Position Enum", "#N/A",
                lambda chart, pos: convert_chart_legend_position(chart, pos)),

    MappingRule("Value Axis Max", "chart_value_axis_max", "Chart", "Value Axis Max",
                "chart.value_axis.maximum_scale", "Float", "#N/A",
                lambda chart, v: convert_chart_value_axis_max(chart, v)),

    MappingRule("Value Axis Min", "chart_value_axis_min", "Chart", "Value Axis Min",
                "chart.value_axis.minimum_scale", "Float", "#N/A",
                lambda chart, v: convert_chart_value_axis_min(chart, v)),

    MappingRule("Series Name", "series_name", "Chart", "Series Name",
                "series.name", "String", "#N/A",
                lambda series, name: convert_series_name(series, name)),

    MappingRule("Series Values", "series_values", "Chart", "Series Values",
                "series.values", "List", "#N/A",
                lambda series, values: convert_series_values(series, values)),

    MappingRule("Series X Values", "series_x_values", "Chart", "Series X Values",
                "series.x_values", "List", "#N/A",
                lambda series, x: convert_series_x_values(series, x)),

    MappingRule("Series Fill Color", "series_fill_color", "Chart", "Series Fill Color",
                "series.format.fill.fore_color.rgb", "RGB tuple", "#N/A",
                lambda series, rgb: convert_series_fill_color(series, rgb)),

    # Notes
    MappingRule("Notes Text", "notes_text", "Notes", "Slide Notes Text",
                "notes_slide.notes_text_frame.text", "String", "#N/A",
                lambda slide, txt: convert_notes_text(slide, txt)),

    # Other
    MappingRule("Hyperlink", "shape_hyperlink", "Other", "Hyperlink",
                "shape.click_action.hyperlink.address", "String", "#N/A",
                lambda shape, url: convert_shape_hyperlink(shape, url)),

    MappingRule("Rotation", "shape_rotation", "Other", "Rotation",
                "shape.rotation", "Degrees", "#N/A",
                lambda shape, deg: convert_shape_rotation(shape, deg)),
]

# -----------------------------
# High-level “apply” helpers
# -----------------------------

def apply_shape_from_bootstrap(shape, class_list: List[str], grid: GridContext,
                               band_height_in: float):
    """Apply geometry/appearance to a shape from Bootstrap-like utilities."""
    p = parse_bootstrap_utils(class_list)

    # Geometry
    if p.offset is not None:
        convert_shape_left_from_offset(shape, grid, p.offset)
    if p.mt is not None:
        convert_shape_top_from_mt(shape, p.mt)
    if p.col is not None:
        convert_shape_width_from_col(shape, grid, p.col)
    if p.h_frac is not None:
        convert_shape_height_from_hfrac(shape, band_height_in, p.h_frac)

    # Appearance
    if p.bg:
        convert_shape_fill_color(shape, p.bg)
    if p.border_zero:
        try: shape.line.fill.background()
        except: pass
    else:
        convert_shape_line_color_default(shape)
        convert_shape_line_width(shape, p.border)
    if p.rounded:
        convert_shape_corner_radius(shape, 0.2)

def apply_paragraph_from_bootstrap(paragraph, text_align_class: Optional[str] = None,
                                   space_before_pt: Optional[float] = None,
                                   space_after_pt: Optional[float] = 6,
                                   line_spacing: Optional[Union[float, int]] = None):
    if text_align_class:
        # "start", "center", "end", "justify"
        convert_paragraph_alignment(paragraph, text_align_class)
    if space_before_pt is not None:
        convert_paragraph_space_before(paragraph, space_before_pt)
    if space_after_pt is not None:
        convert_paragraph_space_after(paragraph, space_after_pt)
    if line_spacing is not None:
        convert_paragraph_line_spacing(paragraph, line_spacing)

def apply_run_from_bootstrap(run, class_list: List[str], size_pt: Optional[float] = None,
                             font_name: Optional[str] = None):
    p = parse_bootstrap_utils(class_list)
    if p.fw_bold:
        convert_run_bold(run, True)
    if p.fst_italic:
        convert_run_italic(run, True)
    if p.text_uppercase:
        convert_run_all_caps(run, True)
    if p.text_color:
        convert_run_color(run, p.text_color)
    if size_pt is not None:
        convert_run_size(run, size_pt)
    if font_name:
        convert_run_name(run, font_name)
