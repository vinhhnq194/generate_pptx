import json
from pptx import Presentation
from pptx.util import Inches
from v5.parsers.layout_parser import parse_layout_only
from v5.layout.grid_solver import solve_layout
from v5.layout.placement_debug import draw_grid, draw_bbox
from v5.renderer.primitives import add_placeholder

def render_layout_only(html_path: str, styles_path: str, template_path: str | None = None):
    ST = json.load(open(styles_path, "r", encoding="utf-8"))
    page, bands = ST["page"], ST["bands"]

    prs = Presentation(template_path) if template_path else Presentation()
    prs.slide_width  = Inches(page["width_in"])
    prs.slide_height = Inches(page["height_in"])
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Parse structure only (rows/cols)
    layout = parse_layout_only(html_path)

    # Solve absolute placements
    placements = solve_layout(layout, page, bands, page["gutter_in"])

    # Optional debug overlays
    if ST.get("debug", {}).get("grid", False):
        draw_grid(slide, page, page["gutter_in"])

    for i, (rect, group) in enumerate(placements, 1):
        label = f"col:{group.span} off:{group.offset}"
        add_placeholder(slide, rect.left, rect.top, rect.width, rect.height, label=label)
        if ST.get("debug", {}).get("bbox", False):
            draw_bbox(slide, rect)

    return prs
