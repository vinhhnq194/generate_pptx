from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

def draw_grid(slide, page: dict, gutter_in: float):
    """Draw 12-column grid as a faint overlay to verify alignment."""
    L, T, R, B = page["margins_in"]
    W = page["width_in"] - L - R
    H = page["height_in"] - T - B
    col_w = (W - 11 * gutter_in) / 12.0

    x = L
    for _ in range(12):
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(T), Inches(col_w), Inches(H))
        rect.fill.solid(); rect.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
        rect.fill.transparency = 0.7
        rect.line.fill.background()
        x += col_w + gutter_in

def draw_bbox(slide, rect, color=RGBColor(0x22, 0x55, 0xAA)):
    """Outline a placement rectangle."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(rect.left), Inches(rect.top), Inches(rect.width), Inches(rect.height))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Inches(0.02)
    return shp
